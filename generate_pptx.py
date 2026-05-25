import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 1. Initialize Presentation & Widescreen (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 2. Color Palette (Premium HSL-mapped tones matching HTML)
    ROYAL_WINE = RGBColor(93, 16, 28)      # #5D101C
    DEEP_WINE = RGBColor(61, 9, 17)        # #3D0911
    METALLIC_GOLD = RGBColor(197, 160, 89) # #C5A059
    SOFT_CHARCOAL = RGBColor(74, 62, 64)   # #4A3E40
    PRISTINE_WHITE = RGBColor(253, 251, 251) # #FDFBFB
    LIGHT_BG = RGBColor(241, 237, 234)     # #f1edea
    TILE_BG = RGBColor(251, 248, 248)      # #FBF8F8
    BORDER_RED = RGBColor(230, 215, 215)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper function to set slide background
    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper function to add elegant left-border slide titles
    def add_slide_header(slide, title_text, dark_mode=False):
        # Vertical accent bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.75), Inches(0.55), Inches(0.06), Inches(0.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = METALLIC_GOLD if dark_mode else ROYAL_WINE
        shape.line.fill.background()
        
        # Title text
        txBox = slide.shapes.add_textbox(Inches(0.95), Inches(0.48), Inches(11.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Playfair Display'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = PRISTINE_WHITE if dark_mode else ROYAL_WINE

    # Helper function to create premium styled tiles
    def add_styled_tile(slide, x, y, cx, cy):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
        shape.fill.solid()
        shape.fill.fore_color.rgb = TILE_BG
        shape.line.color.rgb = BORDER_RED
        shape.line.width = Pt(1.5)
        # Rounded corners adjustment if available in python-pptx (default is clean)
        return shape

    # ==================== SLIDE 1: Problem Statement ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, LIGHT_BG)
    add_slide_header(slide1, "1. The Public Healthcare Crisis")
    
    # 2 Column tiles
    col_width = Inches(5.6)
    col_height = Inches(4.8)
    
    # Column 1
    add_styled_tile(slide1, Inches(0.75), Inches(1.6), col_width, col_height)
    txBox1 = slide1.shapes.add_textbox(Inches(1.15), Inches(1.9), col_width - Inches(0.8), col_height - Inches(0.6))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    
    p_title1 = tf1.paragraphs[0]
    p_title1.text = "⚠️ The Proactive Monitoring Vacuum"
    p_title1.font.name = 'Playfair Display'
    p_title1.font.size = Pt(22)
    p_title1.font.bold = True
    p_title1.font.color.rgb = ROYAL_WINE
    p_title1.space_after = Pt(20)
    
    p_desc1 = tf1.add_paragraph()
    p_desc1.text = "Modern clinical architecture is overwhelmingly reactive. Vulnerable patient groups—including adolescent students with undiagnosed nutrition anomalies (Anaemia/PCOD) and elderly populations—are typically detected only when clinical symptoms escalate to an emergency threshold, incurring high post-crisis management expenses."
    p_desc1.font.name = 'Inter'
    p_desc1.font.size = Pt(15)
    p_desc1.font.color.rgb = SOFT_CHARCOAL
    p_desc1.line_spacing = 1.35
    
    # Column 2
    add_styled_tile(slide1, Inches(6.98), Inches(1.6), col_width, col_height)
    txBox2 = slide1.shapes.add_textbox(Inches(7.38), Inches(1.9), col_width - Inches(0.8), col_height - Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    p_title2 = tf2.paragraphs[0]
    p_title2.text = "🖥️ Systemic & Segregated Data Silos"
    p_title2.font.name = 'Playfair Display'
    p_title2.font.size = Pt(22)
    p_title2.font.bold = True
    p_title2.font.color.rgb = ROYAL_WINE
    p_title2.space_after = Pt(20)
    
    p_desc2 = tf2.add_paragraph()
    p_desc2.text = "Continuous physiological telemetry remains locked or entirely uncaptured. Public schools, rural outreach sites, and geriatric homes fail to integrate physical and digital monitoring networks. This systemic gap blocks medical interventions and prevents district leaders from executing targeted, preventive population care."
    p_desc2.font.name = 'Inter'
    p_desc2.font.size = Pt(15)
    p_desc2.font.color.rgb = SOFT_CHARCOAL
    p_desc2.line_spacing = 1.35

    # ==================== SLIDE 2: Solution & Features ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, LIGHT_BG)
    add_slide_header(slide2, "2. Proactive Dual-Engine Architecture")
    
    # 3 Column tiles
    tile_w = Inches(3.77)
    tile_h = Inches(4.8)
    spacing = Inches(0.26)
    
    modules = [
        {
            "icon": "👶",
            "title": "Module A: Youth Health",
            "text": "Empowers government schools with low-power continuous biometric wearables. Machine learning classification detects early-stage biomarkers of chronic Anaemia and developmental PCOD, allowing stigma-free micro-nutrient support."
        },
        {
            "icon": "❤️",
            "title": "Module B: Geriatric RPM",
            "text": "Safeguards nursing homes and home-care patients via high-fidelity RPM telemetry. Employs low-latency anomaly engines to coordinate automated emergency protocols, bypassing manual wait times during critical midnight hours."
        },
        {
            "icon": "🛡️",
            "title": "Core Technologies",
            "text": "Combines clinical edge computing with private blockchain ledgers. Health data pathways remain strictly immutable and completely consent-controlled, preserving compliance safeguards across sensitive pediatric and geriatric sectors."
        }
    ]
    
    for i, mod in enumerate(modules):
        left_x = Inches(0.75) + i * (tile_w + spacing)
        add_styled_tile(slide2, left_x, Inches(1.6), tile_w, tile_h)
        
        # Textbox inside tile
        tx = slide2.shapes.add_textbox(left_x + Inches(0.3), Inches(1.9), tile_w - Inches(0.6), tile_h - Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        
        # Icon
        p_icon = tf.paragraphs[0]
        p_icon.text = mod["icon"]
        p_icon.font.size = Pt(32)
        p_icon.space_after = Pt(12)
        
        # Title
        p_t = tf.add_paragraph()
        p_t.text = mod["title"]
        p_t.font.name = 'Playfair Display'
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = ROYAL_WINE
        p_t.space_after = Pt(16)
        
        # Description
        p_d = tf.add_paragraph()
        p_d.text = mod["text"]
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(13.5)
        p_d.font.color.rgb = SOFT_CHARCOAL
        p_d.line_spacing = 1.3

    # ==================== SLIDE 3: Business Model ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, LIGHT_BG)
    add_slide_header(slide3, "3. Sustainable Financial Framework")
    
    # Left High Contrast Metric block
    metric_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.6), Inches(5.0), Inches(4.8))
    metric_bg.fill.solid()
    metric_bg.fill.fore_color.rgb = ROYAL_WINE
    metric_bg.line.fill.background()
    
    tx_m = slide3.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(5.0), Inches(3.5))
    tf_m = tx_m.text_frame
    tf_m.word_wrap = True
    
    p_num = tf_m.paragraphs[0]
    p_num.text = "40%"
    p_num.alignment = PP_ALIGN.CENTER
    p_num.font.name = 'Playfair Display'
    p_num.font.size = Pt(100)
    p_num.font.bold = True
    p_num.font.color.rgb = PRISTINE_WHITE
    p_num.space_after = Pt(10)
    
    p_label = tf_m.add_paragraph()
    p_label.text = "LONG-TERM EXPENSE REDUCTIONS"
    p_label.alignment = PP_ALIGN.CENTER
    p_label.font.name = 'Inter'
    p_label.font.size = Pt(14)
    p_label.font.bold = True
    p_label.font.color.rgb = METALLIC_GOLD
    
    # Right Text block
    tx_r = slide3.shapes.add_textbox(Inches(6.25), Inches(1.6), Inches(6.33), Inches(4.8))
    tf_r = tx_r.text_frame
    tf_r.word_wrap = True
    
    p_rt = tf_r.paragraphs[0]
    p_rt.text = "Hardware-Enabled SaaS (HaaS / SaaS)"
    p_rt.font.name = 'Playfair Display'
    p_rt.font.size = Pt(24)
    p_rt.font.bold = True
    p_rt.font.color.rgb = ROYAL_WINE
    p_rt.space_after = Pt(15)
    
    p_rd1 = tf_r.add_paragraph()
    p_rd1.text = "OmniHealth employs a scalable, budget-sensitive monetization model. Public administrations leverage structured lease-to-own plans, while private care networks adopt a predictable subscription tier based on active patient tracking."
    p_rd1.font.name = 'Inter'
    p_rd1.font.size = Pt(14.5)
    p_rd1.font.color.rgb = SOFT_CHARCOAL
    p_rd1.line_spacing = 1.3
    p_rd1.space_after = Pt(15)
    
    p_rd2 = tf_r.add_paragraph()
    p_rd2.text = "Financial Return Equation:"
    p_rd2.font.name = 'Inter'
    p_rd2.font.size = Pt(15)
    p_rd2.font.bold = True
    p_rd2.font.color.rgb = ROYAL_WINE
    p_rd2.space_after = Pt(10)
    
    # Formula container shape
    formula_shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.25), Inches(4.6), Inches(6.33), Inches(1.3))
    formula_shape.fill.solid()
    formula_shape.fill.fore_color.rgb = TILE_BG
    formula_shape.line.color.rgb = ROYAL_WINE
    formula_shape.line.width = Pt(1.5)
    
    # Formula Text
    tx_f = slide3.shapes.add_textbox(Inches(6.35), Inches(4.8), Inches(6.13), Inches(0.9))
    tf_f = tx_f.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "ROI = [ Savings (Avoided Emergency) - Cost (Ops) ] / Cost (Ops) x 100%"
    p_f.alignment = PP_ALIGN.CENTER
    p_f.font.name = 'Inter'
    p_f.font.size = Pt(13)
    p_f.font.bold = True
    p_f.font.color.rgb = ROYAL_WINE

    # ==================== SLIDE 4: Experience ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, LIGHT_BG)
    add_slide_header(slide4, "4. Seamless Care Integration Journey")
    
    # Horizontal timeline line
    tl_line = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.9), Inches(11.83), Inches(0.04))
    tl_line.fill.solid()
    tl_line.fill.fore_color.rgb = ROYAL_WINE
    tl_line.line.fill.background()
    
    steps = [
        {
            "step": "1. Wear & Align",
            "desc": "Patients receive non-invasive, skin-safe biometric stress bands or medical tracker kits requiring minimal setup.",
            "is_above": True
        },
        {
            "step": "2. Passive Syncing",
            "desc": "Vitals stream silently to secure bedside or edge gateways without requiring patient interaction.",
            "is_above": False
        },
        {
            "step": "3. AI Stratification",
            "desc": "Proprietary cloud engines compute anomalies and generate priority-stratified clinical flags.",
            "is_above": True
        },
        {
            "step": "4. Command Response",
            "desc": "Administrators inspect centralized regional heatmaps to deploy targeted healthcare camps immediately.",
            "is_above": False
        }
    ]
    
    step_w = Inches(2.74)
    step_h = Inches(1.9)
    step_spacing = Inches(0.29)
    
    for idx, st in enumerate(steps):
        left_x = Inches(0.75) + idx * (step_w + step_spacing)
        
        # Compute Y position based on alternate timeline structure
        if st["is_above"]:
            top_y = Inches(1.7)
        else:
            top_y = Inches(4.35)
            
        # Draw step box
        add_styled_tile(slide4, left_x, top_y, step_w, step_h)
        
        # Step text
        tx = slide4.shapes.add_textbox(left_x + Inches(0.15), top_y + Inches(0.15), step_w - Inches(0.3), step_h - Inches(0.3))
        tf = tx.text_frame
        tf.word_wrap = True
        
        p_st = tf.paragraphs[0]
        p_st.text = st["step"]
        p_st.font.name = 'Playfair Display'
        p_st.font.size = Pt(17)
        p_st.font.bold = True
        p_st.font.color.rgb = ROYAL_WINE
        p_st.space_after = Pt(8)
        
        p_sd = tf.add_paragraph()
        p_sd.text = st["desc"]
        p_sd.font.name = 'Inter'
        p_sd.font.size = Pt(11.5)
        p_sd.font.color.rgb = SOFT_CHARCOAL
        p_sd.line_spacing = 1.25
        
        # Timeline node dot
        dot_cx = left_x + (step_w / 2)
        dot_size = Inches(0.18)
        dot = slide4.shapes.add_shape(MSO_SHAPE.OVAL, dot_cx - (dot_size / 2), Inches(3.83), dot_size, dot_size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRISTINE_WHITE
        dot.line.color.rgb = METALLIC_GOLD
        dot.line.width = Pt(3)

    # ==================== SLIDE 5: Conclusion (Dark Theme) ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, DEEP_WINE)
    
    # 1. Badge "Platform Paradigm Shift"
    badge = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.91), Inches(1.4), Inches(3.5), Inches(0.45))
    badge.fill.solid()
    # Light gold semi-transparent look simulated via solid light color
    badge.fill.fore_color.rgb = RGBColor(90, 20, 30)
    badge.line.color.rgb = METALLIC_GOLD
    badge.line.width = Pt(1.5)
    
    tx_b = slide5.shapes.add_textbox(Inches(4.91), Inches(1.46), Inches(3.5), Inches(0.4))
    tf_b = tx_b.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "PLATFORM PARADIGM SHIFT"
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.name = 'Inter'
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = METALLIC_GOLD
    
    # 2. Main Title
    tx_t = slide5.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.2))
    tf_t = tx_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "The Proactive Healthcare Horizon"
    p_t.alignment = PP_ALIGN.CENTER
    p_t.font.name = 'Playfair Display'
    p_t.font.size = Pt(54)
    p_t.font.bold = True
    p_t.font.color.rgb = PRISTINE_WHITE
    
    # 3. Subtitle / Quote
    tx_q = slide5.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.33), Inches(1.0))
    tf_q = tx_q.text_frame
    tf_q.word_wrap = True
    p_q = tf_q.paragraphs[0]
    p_q.text = '"Transforming public administration from reactive crisis response to proactive preventative containment."'
    p_q.alignment = PP_ALIGN.CENTER
    p_q.font.name = 'Inter'
    p_q.font.size = Pt(18)
    p_q.font.italic = True
    p_q.font.color.rgb = RGBColor(230, 220, 220)
    
    # 4. Horizontal line
    f_line = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.66), Inches(4.9), Inches(6.0), Inches(0.02))
    f_line.fill.solid()
    f_line.fill.fore_color.rgb = RGBColor(120, 40, 50)
    f_line.line.fill.background()

    # 5. Contact Info
    tx_c = slide5.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.33), Inches(0.8))
    tf_c = tx_c.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.text = "🌐  omnihealth-systems.org      |      ✉️  admin@omnihealth-systems.org"
    p_c.alignment = PP_ALIGN.CENTER
    p_c.font.name = 'Inter'
    p_c.font.size = Pt(15)
    p_c.font.bold = True
    p_c.font.color.rgb = METALLIC_GOLD

    # 6. Save presentation to Downloads folder
    downloads_path = os.path.join(os.path.expanduser('~'), 'Downloads')
    file_path = os.path.join(downloads_path, "OmniHealth_Proposal_Presentation.pptx")
    
    prs.save(file_path)
    print(f"Presentation saved successfully to: {file_path}")
    return file_path

if __name__ == "__main__":
    create_presentation()
